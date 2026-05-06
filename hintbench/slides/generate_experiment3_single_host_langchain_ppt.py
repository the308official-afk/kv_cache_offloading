from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT = Path('hintbench/slides/HintBench_Experiment3_Single_Host_LangChain_Client_Path_Run.pptx')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(248, 250, 252)
NAVY = RGBColor(17, 24, 39)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
AMBER = RGBColor(217, 119, 6)
ROSE = RGBColor(225, 29, 72)
GRAY = RGBColor(107, 114, 128)
LIGHT = RGBColor(229, 231, 235)
WHITE = RGBColor(255, 255, 255)
SOFT_BLUE = RGBColor(219, 234, 254)
SOFT_GREEN = RGBColor(209, 250, 229)
SOFT_ORANGE = RGBColor(254, 243, 199)
SOFT_PINK = RGBColor(255, 228, 230)
SOFT_GRAY = RGBColor(243, 244, 246)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.62), Inches(1.05), Inches(12.0), Inches(0.5))
        tf2 = st.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(13)
        r2.font.color.rgb = GRAY


def add_bullets(slide, items, left=0.9, top=1.5, width=11.5, height=5.3, font_size=20, color=NAVY):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(10)


def add_box(slide, left, top, width, height, text, fill_color, line_color=WHITE, font_size=15, bold=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = NAVY
    return shape


def add_arrow(slide, left, top, width=0.5):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(left), Inches(top), Inches(width), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.color.rgb = BLUE
    return shape


def add_placeholder_box(slide, title, note):
    add_title(slide, title, note)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.8), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = SOFT_GRAY
    box.line.color.rgb = LIGHT
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = 'Paste REQUEST_LOG_MODE = "full" output here after a sample run.\n\nSuggested content:\n- raw log excerpt\n- short interpretation\n- anything unexpected\n\nPlaceholder:'
    r.font.size = Pt(18)
    r.font.color.rgb = GRAY
    inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(3.15), Inches(11.3), Inches(2.65))
    inner.fill.solid()
    inner.fill.fore_color.rgb = WHITE
    inner.line.color.rgb = LIGHT
    tf2 = inner.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = '[ SAMPLE LOG GOES HERE ]'
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = LIGHT
    p2.alignment = PP_ALIGN.CENTER
    tf2.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'HintBench Experiment 3', 'Single-Host LangChain Client-Path Run')
add_box(slide, 0.8, 1.5, 11.7, 1.0,
        'Best one-line summary: HintBench Experiment 3 is a single-host, 4-stage LangChain client-path benchmark that sends hint-carrying requests directly to a local Dynamo frontend and one local SGLang worker, with full client-side checkpoint logging.',
        SOFT_BLUE, LIGHT, 20, True)
add_bullets(slide, [
    'Scope: only Experiment 3, with LangChain present and no shim.',
    'Pipeline type: request generator -> LangChain -> Dynamo frontend -> SGLang worker.',
    'Run target: localhost frontend on port 8000 via run_dynamo_single_host.sh.',
], top=3.0, font_size=22)

# Slide 2 setup/components
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'Setup Components', 'All components used in Experiment 3')
add_box(slide, 0.7, 1.5, 2.2, 0.85, 'Experiment YAML\nbaseline_round_robin_langchain.yaml', SOFT_BLUE)
add_box(slide, 3.1, 1.5, 1.7, 0.85, 'run_experiment.py', SOFT_GREEN)
add_box(slide, 5.0, 1.5, 1.7, 0.85, 'shared_prefix.py\nrequest generator', SOFT_ORANGE)
add_box(slide, 6.9, 1.5, 1.8, 0.85, 'langchain_loadgen.py', SOFT_PINK)
add_box(slide, 8.9, 1.5, 1.5, 0.85, 'ChatOpenAI', SOFT_BLUE)
add_box(slide, 10.6, 1.5, 2.0, 0.85, 'Dynamo frontend\nlocalhost:8000', SOFT_GREEN)
for x in [2.95,4.85,6.75,8.75,10.45]:
    add_arrow(slide, x, 1.7)
add_box(slide, 3.0, 3.2, 3.3, 1.1, 'run_dynamo_single_host.sh\nstarts and verifies the single-host stack', WHITE, LIGHT, 18, True)
add_box(slide, 6.7, 3.2, 2.8, 1.1, 'Single local\nSGLang worker', SOFT_ORANGE, LIGHT, 18, True)
add_box(slide, 9.9, 3.2, 2.6, 1.1, 'HintBench outputs\nmetadata.json / results.jsonl / summary.json', SOFT_PINK, LIGHT, 16, True)
add_bullets(slide, [
    'No live_hint_router.py in this experiment.',
    'Frontend URL must be http://127.0.0.1:8000/v1/chat/completions.',
    'Default worker flags: --enable-cache-report --enable-priority-scheduling --radix-eviction-policy lru.',
], top=4.8, font_size=18)

# Slide 3 processing flow
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'Processing Flow', 'How one request moves through the Experiment 3 path')
flow_items = [
    '1. Experiment config selects client_backend = langchain and shared_prefix workload generation.',
    '2. HintBench generates raw request rows with messages, request_id, prompt_id, and hint_payload.',
    '3. LangChain path reads each request row and logs it before transformation.',
    '4. Raw messages are converted into LangChain message objects: SystemMessage, HumanMessage, AIMessage.',
    '5. hint_payload is wrapped into nvext.agent_hints using LangChain extra_body.',
    '6. ChatOpenAI sends the OpenAI-style request to the local Dynamo frontend.',
    '7. Dynamo frontend routes the request to the single local SGLang worker for inference.',
    '8. HintBench writes normal run outputs and prints checkpoint logs to stdout.',
]
add_bullets(slide, flow_items, top=1.5, font_size=18)

# Slide 4 setup/run commands
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'Run Commands', 'Exact commands for the single-host LangChain client-path run')
cmd_text = (
    'Start / verify the single-host stack\n'
    'cd ~/kv_cache_offloading\n'
    './run_dynamo_single_host.sh start\n'
    './run_dynamo_single_host.sh status\n'
    './run_dynamo_single_host.sh test\n\n'
    'Install LangChain if needed\n'
    'python3 -m pip install -U langchain-openai langchain-core\n\n'
    'Run the experiment\n'
    'python3 hintbench/run_experiment.py \\\n'
    '  --config hintbench/experiments/baseline_round_robin_langchain.yaml \\\n'
    '  --frontend-url http://127.0.0.1:8000/v1/chat/completions'
)
box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.45), Inches(12.0), Inches(4.8))
box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.color.rgb = NAVY
textf = box.text_frame; textf.word_wrap = True
p = textf.paragraphs[0]
r = p.add_run(); r.text = cmd_text; r.font.name = 'Courier New'; r.font.size = Pt(17); r.font.color.rgb = WHITE
add_bullets(slide, [
    'Inspect the latest run with summary.json, results.jsonl, and metadata.json.',
    'Use REQUEST_LOG_MODE = "full" in hintbench/constants.py for maximum checkpoint detail.',
], top=6.45, font_size=17)

# Slide 5 capabilities
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'Current Capabilities', 'What this experiment can do today')
add_bullets(slide, [
    'Generate synthetic shared-prefix multi-turn requests with stable hint payloads.',
    'Send requests through LangChain instead of the default async client.',
    'Show the request before and after LangChain transformation.',
    'Show where hints are injected under nvext.agent_hints.',
    'Send all requests through a local single-host Dynamo + SGLang stack.',
    'Produce standard HintBench result artifacts plus detailed client-path checkpoint logs.',
    'Support phase-aware debugging of the client path without involving the live shim.',
], top=1.55, font_size=19)

# Slide 6 hints
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'Hints Currently Supported', 'Hints relevant to Experiment 3')
add_box(slide, 0.8, 1.45, 2.0, 0.7, 'NVIDIA hint used', SOFT_BLUE, LIGHT, 18, True)
add_box(slide, 3.0, 1.45, 9.5, 0.7, 'priority: marks urgency / importance and is the main NVIDIA-documented hint used end to end here.', WHITE, LIGHT, 17)
customs = [
    'reuse_likelihood: likelihood the request benefits from existing KV cache or repeated prefixes.',
    'agent_phase: workflow phase label; current experiment defaults to execution.',
    'latency_sensitivity: how strongly lower wait time should be preferred.',
    'program_id: which workflow generated the request; here, hintbench.shared_prefix.',
    'context_type: kind of context being used; here, multi_turn_shared_prefix.',
    'expected_output_tokens: estimated response length.',
]
add_bullets(slide, customs, top=2.5, font_size=17)
add_bullets(slide, [
    'Not currently supported end to end in this setup: osl, speculative_prefill, session_control.',
], top=6.2, font_size=17, color=ROSE)

# Slide 7 config snapshot
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'Experiment Snapshot', 'Key settings from baseline_round_robin_langchain.yaml')
config_text = (
    'name: baseline_round_robin_langchain\n'
    'router_mode: round-robin\n'
    'model: Qwen/Qwen2.5-0.5B\n'
    'workload: shared_prefix\n'
    'client_backend: langchain\n'
    'concurrency: 4\n'
    'num_conversations: 4\n'
    'turns_per_conversation: 3\n'
    'max_tokens: 128\n'
    'temperature: 0.0\n'
    'request_timeout_s: 120\n'
    'results_timezone: America/Chicago\n'
    'shared_prefix_group: group-a'
)
box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.5), Inches(5.2), Inches(4.9))
box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.color.rgb = NAVY
textf = box.text_frame
p = textf.paragraphs[0]
r = p.add_run(); r.text = config_text; r.font.name='Courier New'; r.font.size = Pt(18); r.font.color.rgb = WHITE
add_box(slide, 6.5, 1.7, 5.8, 1.8, 'Default hint_defaults\npriority=5, reuse_likelihood=0.9, agent_phase=execution, latency_sensitivity=0.7, program_id=hintbench.shared_prefix, context_type=multi_turn_shared_prefix, expected_output_tokens=128', SOFT_ORANGE, LIGHT, 16, True)
add_box(slide, 6.5, 4.0, 5.8, 1.4, 'Log settings\nREQUEST_LOG_MODE = "full"\nCONVERTED_MESSAGE_LOG_MODE = "full"\nHINT_INJECTION_LOG_MODE = "full"\nREQUEST_DISPATCH_LOG_MODE = "full"', SOFT_GREEN, LIGHT, 16, True)

# Slide 8 outputs and observability
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'Outputs and Observability', 'What you get after a run')
add_bullets(slide, [
    'Run directory under hintbench/results/baseline_round_robin_langchain_<timestamp>/',
    'metadata.json: experiment settings and run metadata.',
    'workload.jsonl: generated request rows.',
    'results.jsonl: per-request results.',
    'summary.json: aggregate run summary.',
    'stdout checkpoint logs from langchain_loadgen.py for request receipt, message conversion, hint injection, and dispatch.',
    'Single-host service logs via ./run_dynamo_single_host.sh logs or logs -f.',
], top=1.55, font_size=18)

# Slide 9 checkpoint map
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'Checkpoint Map', 'Where the full client-path logs appear in processing')
add_box(slide, 0.8, 1.55, 5.8, 0.9, '1. Request received\nGenerated HintBench request rows are printed before LangChain converts messages or sends the request onward.', SOFT_BLUE, LIGHT, 17, True)
add_box(slide, 6.8, 1.55, 5.8, 0.9, '2. Messages converted\nLangChain message objects are logged immediately after conversion.', SOFT_GREEN, LIGHT, 17, True)
add_box(slide, 0.8, 3.0, 5.8, 0.9, '3. Hints injected\nHints are injected into nvext.agent_hints via LangChain extra_body.', SOFT_ORANGE, LIGHT, 17, True)
add_box(slide, 6.8, 3.0, 5.8, 0.9, '4. Request dispatched\nLangChain sends the request to the frontend here.', SOFT_PINK, LIGHT, 17, True)
add_bullets(slide, [
    'These placeholders are intended for REQUEST_LOG_MODE = "full" after a sample run.',
    'Each checkpoint can later be annotated with what changed from the previous stage.',
], top=4.6, font_size=18)

# Slides 10-13 placeholders
add_placeholder_box(prs.slides.add_slide(prs.slide_layouts[6]), 'Checkpoint Placeholder 1', 'Generated HintBench request rows are printed here before LangChain converts messages or sends the request onward.')
add_placeholder_box(prs.slides.add_slide(prs.slide_layouts[6]), 'Checkpoint Placeholder 2', 'LangChain message objects are logged here immediately after conversion so they can be compared against the incoming raw messages.')
add_placeholder_box(prs.slides.add_slide(prs.slide_layouts[6]), 'Checkpoint Placeholder 3', 'Hints are injected into the outgoing request here under nvext.agent_hints via LangChain extra_body.')
add_placeholder_box(prs.slides.add_slide(prs.slide_layouts[6]), 'Checkpoint Placeholder 4', 'LangChain sends the request to the frontend here.')

# Slide 14 important notes
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, 'Important Notes', 'Experiment-specific boundaries and caveats')
add_bullets(slide, [
    'This deck is only for Experiment 3: Single-Host LangChain Client-Path Run.',
    'Do not start live_hint_router.py for this experiment.',
    'Use port 8000, not 8100.',
    'This experiment is best for client-path visibility, not for multi-worker routing comparisons.',
    'The LangChain path does not currently capture Dynamo-specific fields like ttft_ms, kv_hit_rate, cached_tokens, or worker_id.',
    'Single-host is best for iteration and debugging, but it is not a substitute for the full multi-worker setup.',
], top=1.55, font_size=19)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
