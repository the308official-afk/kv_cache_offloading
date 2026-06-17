#!/usr/bin/env python3
from __future__ import annotations

import argparse
import contextlib
import http.server
import socketserver
import tempfile
import threading
from pathlib import Path
from urllib.parse import quote

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Convert a local HTML slide deck into a PPTX by rendering each slide "
            "element as an image and placing one image per PowerPoint slide."
        )
    )
    parser.add_argument("html_path", help="Path to the HTML slide deck")
    parser.add_argument(
        "--output",
        default=None,
        help="Output PPTX path (defaults to the HTML path with a .pptx suffix)",
    )
    parser.add_argument(
        "--selector",
        default=".slide",
        help="CSS selector that identifies each slide element (default: .slide)",
    )
    parser.add_argument(
        "--timeout-ms",
        type=int,
        default=30000,
        help="Navigation and asset wait timeout in milliseconds",
    )
    parser.add_argument(
        "--scale",
        type=float,
        default=2.0,
        help="Device scale factor for slide screenshots (default: 2.0)",
    )
    return parser.parse_args()


class ThreadedTCPServer(socketserver.ThreadingMixIn, socketserver.TCPServer):
    allow_reuse_address = True
    daemon_threads = True


@contextlib.contextmanager
def serve_directory(directory: Path):
    handler = lambda *args, **kwargs: http.server.SimpleHTTPRequestHandler(  # noqa: E731
        *args, directory=str(directory), **kwargs
    )
    with ThreadedTCPServer(("127.0.0.1", 0), handler) as server:
        thread = threading.Thread(target=server.serve_forever, daemon=True)
        thread.start()
        try:
            yield server
        finally:
            server.shutdown()
            thread.join(timeout=2)


def wait_for_assets(page, timeout_ms: int) -> None:
    page.wait_for_load_state("networkidle", timeout=timeout_ms)
    page.wait_for_function(
        "() => document.fonts ? document.fonts.status === 'loaded' : true",
        timeout=timeout_ms,
    )
    page.wait_for_function(
        """
        () => Array.from(document.images).every((img) => img.complete)
        """,
        timeout=timeout_ms,
    )


def detect_slide_size(page, selector: str) -> tuple[float, float]:
    box = page.locator(selector).first.bounding_box()
    if not box:
        raise SystemExit(f"Could not determine bounding box for selector: {selector}")
    return box["width"], box["height"]


def ppt_dimensions(width_px: float, height_px: float) -> tuple[float, float]:
    # Keep a standard wide-slide width and derive height from the detected aspect ratio.
    slide_width_in = 13.333333
    slide_height_in = slide_width_in * (height_px / width_px)
    return slide_width_in, slide_height_in


def build_pptx(image_paths: list[Path], width_px: float, height_px: float, output_path: Path) -> None:
    prs = Presentation()
    slide_width_in, slide_height_in = ppt_dimensions(width_px, height_px)
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)
    blank_layout = prs.slide_layouts[6]

    for image_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    # python-pptx creates one empty default slide on initialization.
    xml_slides = prs.slides._sldIdLst  # type: ignore[attr-defined]
    del xml_slides[0]
    prs.save(output_path)


def render_slides(html_path: Path, selector: str, timeout_ms: int, scale: float) -> tuple[list[Path], float, float]:
    html_path = html_path.resolve()
    if not html_path.exists():
        raise SystemExit(f"HTML file not found: {html_path}")

    with tempfile.TemporaryDirectory(prefix="html-slides-to-pptx-") as temp_dir_name:
        temp_dir = Path(temp_dir_name)
        screenshots_dir = temp_dir / "screens"
        screenshots_dir.mkdir(parents=True, exist_ok=True)

        with serve_directory(html_path.parent) as server, sync_playwright() as playwright:
            browser = playwright.chromium.launch(headless=True)
            page = browser.new_page(
                viewport={"width": 1600, "height": 1200},
                device_scale_factor=scale,
            )

            relative_url = quote(html_path.name)
            page.goto(f"http://127.0.0.1:{server.server_address[1]}/{relative_url}", wait_until="domcontentloaded")
            wait_for_assets(page, timeout_ms)

            slides = page.locator(selector)
            count = slides.count()
            if count == 0:
                browser.close()
                raise SystemExit(f"No slides matched selector {selector!r} in {html_path}")

            width_px, height_px = detect_slide_size(page, selector)
            image_paths: list[Path] = []
            for index in range(count):
                path = screenshots_dir / f"slide-{index + 1:02d}.png"
                slides.nth(index).screenshot(path=str(path))
                image_paths.append(path)

            browser.close()

            persisted_dir = html_path.parent / f"{html_path.stem}-pptx-assets"
            if persisted_dir.exists():
                for old in persisted_dir.glob("slide-*.png"):
                    old.unlink()
            persisted_dir.mkdir(exist_ok=True)

            persisted_paths: list[Path] = []
            for image_path in image_paths:
                target = persisted_dir / image_path.name
                target.write_bytes(image_path.read_bytes())
                persisted_paths.append(target)

            return persisted_paths, width_px, height_px


def main() -> int:
    args = parse_args()
    html_path = Path(args.html_path).resolve()
    output_path = (
        Path(args.output).resolve()
        if args.output
        else html_path.with_suffix(".pptx")
    )

    image_paths, width_px, height_px = render_slides(
        html_path=html_path,
        selector=args.selector,
        timeout_ms=args.timeout_ms,
        scale=args.scale,
    )
    build_pptx(image_paths, width_px, height_px, output_path)

    print(f"Generated PPTX: {output_path}")
    print(f"Slides captured: {len(image_paths)}")
    print(f"Selector: {args.selector}")
    print(f"Rendered slide size: {width_px:.0f}x{height_px:.0f}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
