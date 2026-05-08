from __future__ import annotations

import json
import textwrap
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
RUN_DIR = ROOT / "agentbench/results/sample__stronger_behavior__001_20260506_155649"
OUT = ROOT / "agentbench/slides/AgentBench_Experiment1_Upstream_Deploy_Coding_Agent_Run.pptx"


BG = RGBColor(248, 250, 252)
NAVY = RGBColor(17, 24, 39)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
AMBER = RGBColor(217, 119, 6)
ROSE = RGBColor(225, 29, 72)
GRAY = RGBColor(107, 114, 128)
LIGHT = RGBColor(229, 231, 235)
WHITE = RGBColor(255, 255, 255)
SOFT_BLUE = RGBColor(219, 234, 254)
SOFT_GREEN = RGBColor(209, 250, 229)
SOFT_ORANGE = RGBColor(254, 243, 199)
SOFT_PINK = RGBColor(255, 228, 230)
SOFT_GRAY = RGBColor(243, 244, 246)


def load_json(path: Path):
    return json.loads(path.read_text())


def shorten(value: str, width: int = 160) -> str:
    single = " ".join(value.split())
    return textwrap.shorten(single, width=width, placeholder="...")


def nonempty_lines(text: str) -> list[str]:
    return [line for line in text.splitlines() if line.strip()]


def excerpt_after(text: str, marker: str, *, max_lines: int = 6) -> list[str]:
    lines = nonempty_lines(text)
    for idx, line in enumerate(lines):
        if marker in line:
            return lines[idx : idx + max_lines]
    return lines[:max_lines]


def wrap_excerpt(text: str, *, width: int = 96, max_lines: int = 7) -> list[str]:
    cleaned = " ".join(str(text).split())
    wrapped = textwrap.wrap(cleaned, width=width)
    if len(wrapped) > max_lines:
        wrapped = wrapped[: max_lines - 1] + [wrapped[max_lines - 1] + " ..."]
    return wrapped


def extract_saved_ai_text(saved_response: object) -> str:
    if isinstance(saved_response, dict):
        messages = saved_response.get("messages")
        if isinstance(messages, list) and messages:
            for msg in reversed(messages):
                if isinstance(msg, dict) and msg.get("type") == "ai":
                    return str(msg.get("content", ""))
        content = saved_response.get("content")
        if content is not None:
            return str(content)
    return str(saved_response)


def truncate_for_slide(value: Any, *, max_string: int = 240) -> Any:
    if isinstance(value, dict):
        return {k: truncate_for_slide(v, max_string=max_string) for k, v in value.items()}
    if isinstance(value, list):
        return [truncate_for_slide(v, max_string=max_string) for v in value]
    if isinstance(value, str):
        collapsed = value.replace("\n", "\\n")
        if len(collapsed) > max_string:
            return collapsed[: max_string - 16] + "... [truncated]"
        return collapsed
    return value


def raw_json_lines(payload: dict[str, Any], *, max_string: int = 240, max_lines: int = 30) -> list[str]:
    dumped = json.dumps(truncate_for_slide(payload, max_string=max_string), indent=2)
    lines = dumped.splitlines()
    if len(lines) > max_lines:
        lines = lines[: max_lines - 1] + ["  ..."]
    return lines


checkpoints = load_json(RUN_DIR / "checkpoints.json")
result = load_json(RUN_DIR / "result.json")
plan = load_json(RUN_DIR / "plan.json")
final_summary = (RUN_DIR / "final_summary.txt").read_text()

checkpoint_1 = next(item for item in checkpoints if item["check_point"].startswith("1."))
checkpoint_3 = next(item for item in checkpoints if item["check_point"].startswith("3."))
checkpoint_4s = [item for item in checkpoints if item["check_point"].startswith("4.")]
checkpoint_5 = next(item for item in checkpoints if item["check_point"].startswith("5."))
checkpoint_4_focus = checkpoint_4s[-1]

task = checkpoint_1["task"]
step_titles = [step["title"] for step in plan["steps"]]
useful_steps = [title for title in step_titles if "dry-run" in title.lower() or "dispatch" in title.lower() or "temporary" in title.lower()]
generic_steps = [title for title in step_titles if title not in useful_steps]
cp3_excerpt = excerpt_after(checkpoint_3["prompt"], "Return only valid JSON", max_lines=8)
cp3_guidance_excerpt = excerpt_after(checkpoint_3["prompt"], "Prefer task-specific steps", max_lines=6)
cp4_plan_excerpt = excerpt_after(checkpoint_4_focus["prompt"], "Approved decomposition plan:", max_lines=6)
cp4_current_step_excerpt = excerpt_after(checkpoint_4_focus["prompt"], "Current step:", max_lines=6)
cp4_completed_summary_excerpt = excerpt_after(checkpoint_4_focus["prompt"], "Completed step summaries so far:", max_lines=4)
cp5_step_results_excerpt = excerpt_after(checkpoint_5["prompt"], "Step results:", max_lines=5)
cp5_final_job_excerpt = excerpt_after(checkpoint_5["prompt"], "Produce a final summary with these sections:", max_lines=6)
step4_result = result["step_results"][3]
step4_output_excerpt = wrap_excerpt(extract_saved_ai_text(step4_result["response"]), width=94, max_lines=7)
final_output_excerpt = wrap_excerpt(result["result"].get("response_text", final_summary), width=94, max_lines=8)
plan_output_lines = [f"{step['step_id']}. {step['title']}" for step in plan["steps"]]


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.62), Inches(1.05), Inches(12.0), Inches(0.5))
        tf2 = st.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(13)
        r2.font.color.rgb = GRAY


def add_bullets(slide, items, left=0.9, top=1.5, width=11.5, height=5.3, font_size=20, color=NAVY):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_box(slide, left, top, width, height, text, fill_color, line_color=WHITE, font_size=15, bold=False, color=NAVY):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


def add_text_panel(slide, left, top, width, height, title, body, fill=WHITE, title_color=NAVY):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = LIGHT
    tf = panel.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = title_color

    for idx, line in enumerate(body.split("\n")):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY
        p.space_after = Pt(4)
        if idx == 0:
            p.space_before = Pt(6)
    return panel


def add_sections_text(
    slide,
    sections,
    *,
    left=0.85,
    top=1.55,
    width=11.7,
    height=5.6,
    section_title_size=18,
    body_size=15,
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    first = True
    for title, body_lines in sections:
        p_title = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p_title.space_before = Pt(0 if p_title is tf.paragraphs[0] else 10)
        p_title.space_after = Pt(2)
        r_title = p_title.add_run()
        r_title.text = title
        r_title.font.size = Pt(section_title_size)
        r_title.font.bold = True
        r_title.font.color.rgb = BLUE

        for line in body_lines:
            p_body = tf.add_paragraph()
            p_body.text = line
            p_body.font.size = Pt(body_size)
            p_body.font.color.rgb = NAVY
            p_body.space_after = Pt(3)


def add_preformatted_text(
    slide,
    lines,
    *,
    left=0.85,
    top=1.65,
    width=11.6,
    height=5.5,
    font_size=11,
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "\n".join(lines)
    r.font.name = "Courier New"
    r.font.size = Pt(font_size)
    r.font.color.rgb = NAVY


def add_arrow(slide, left, top, width=0.5):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(left), Inches(top), Inches(width), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.color.rgb = BLUE
    return shape


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "AgentBench Experiment 1", "Upstream Deploy-Coding-Agent Variant Run")
add_box(
    slide,
    0.8,
    1.5,
    11.7,
    1.0,
    "Best one-line summary: complex tasks from sources like SWE-bench Pro are broken into several steps by the Deep Agents app before the resulting planning, execution, and synthesis requests are sent through the local Dynamo frontend.",
    SOFT_BLUE,
    LIGHT,
    20,
    True,
)
add_bullets(
    slide,
    [
        "Scope: only Experiment 1, using the upstream deploy-coding-agent variant as the core testbed.",
        "Evidence source: real run sample__stronger_behavior__001_20260506_155649.",
        "This deck focuses on how the payload evolves across checkpoints 1, 3, 4, and 5.",
    ],
    top=3.0,
    font_size=22,
)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Setup Components", "All components used in the upstream variant run")
add_box(slide, 0.6, 1.55, 1.8, 0.85, "Task source\nsample JSON or SWE-bench Pro", SOFT_BLUE)
add_box(slide, 2.7, 1.55, 1.7, 0.85, "AgentBench runner\nsingle-host wrapper", SOFT_GREEN)
add_box(slide, 4.7, 1.55, 1.8, 0.85, "Deep Agents app\nactive harness", SOFT_ORANGE)
add_box(slide, 6.8, 1.55, 2.0, 0.85, "Upstream deploy-coding-agent\ninstructions + skills", SOFT_PINK)
add_box(slide, 9.1, 1.55, 1.6, 0.85, "Dynamo frontend\nlocalhost:8000", SOFT_BLUE)
add_box(slide, 11.0, 1.55, 1.7, 0.85, "Local SGLang\nworker", SOFT_GREEN)
for x in [2.45, 4.45, 6.55, 8.85, 10.75]:
    add_arrow(slide, x, 1.75)
add_box(slide, 1.0, 3.2, 3.2, 1.0, "Checkpoint logging\n1 / 3 / 4 / 5 to stdout + checkpoints.json", WHITE, LIGHT, 17, True)
add_box(slide, 4.6, 3.2, 3.2, 1.0, "Saved artifacts\nresult.json / plan.json / step_results.json / final_summary.txt", WHITE, LIGHT, 17, True)
add_box(slide, 8.2, 3.2, 3.8, 1.0, "Runtime note\nthis run still resolved deepagents from python_environment, not the cloned repo path", SOFT_ORANGE, LIGHT, 16, True)
add_bullets(
    slide,
    [
        "The app variant is upstream_deploy_coding_agent.",
        "The stronger-behavior helper used an easier sample task and a 4-step budget.",
    ],
    top=4.85,
    font_size=18,
)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Processing Flow", "How one experiment run moves through the stack")
add_bullets(
    slide,
    [
        "1. AgentBench loads one task from a local sample JSON or SWE-bench Pro.",
        "2. Checkpoint 1 logs the normalized task before Deep Agents planning starts.",
        "3. The Deep Agents app sends a planning prompt to Dynamo and logs Checkpoint 3.",
        "4. The resulting plan becomes one or more step execution requests and each is logged at Checkpoint 4.",
        "5. After the step results come back, the app sends one final synthesis request and logs Checkpoint 5.",
        "6. AgentBench saves result.json, checkpoints.json, plan.json, step_results.json, and final_summary.txt.",
    ],
    top=1.55,
    font_size=19,
)

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "How The Pieces Connect", "Single-GPU setup in plain language")
add_bullets(
    slide,
    [
        "SWE-bench Pro is the task source. It provides one hard software-engineering task.",
        "The AgentBench runner loads that task and starts one local run.",
        "The Deep Agents app is the agent harness. It turns the task into planning, step-execution, and synthesis requests.",
        "Those requests are sent to the local Dynamo frontend at http://127.0.0.1:8000/v1/chat/completions.",
        "Dynamo forwards the requests to the single local SGLang worker, which runs the model and returns the responses.",
        "AgentBench then saves the run artifacts such as result.json, plan.json, step_results.json, final_summary.txt, and checkpoints.json.",
    ],
    top=1.55,
    font_size=19,
)

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Why We Still Wrote Scripts", "Why the upstream Deep Agents repo was not enough by itself")
add_bullets(
    slide,
    [
        "The Deep Agents GitHub repo gives us the agent framework and example app patterns, but not this exact single-GPU benchmark pipeline.",
        "Our scripts are the adapter layer that connects: SWE-bench task loading, local sample tasks, Deep Agents workflow, Dynamo frontend, SGLang, and saved run artifacts.",
        "The upstream repo does not automatically provide our checkpoint logging, our results directory structure, or our exact AgentBench artifact files.",
        "So the GitHub repo provides the engine and examples, while our scripts provide the local wrapper, experiment flow, and observability for this setup.",
    ],
    top=1.55,
    font_size=19,
)

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "What Upstream Does Not Know By Default", "Added verbatim from the AgentBench setup explanation")
add_bullets(
    slide,
    [
        "1. It does not know about your SWE-bench ingestion flow",
        "Upstream Deep Agents does not automatically do:",
        "load one task from ScaleAI/SWE-bench_Pro",
        "or load your local sample JSON",
        "or save run artifacts in your repo’s format",
        "That part is specific to your benchmark harness.",
        "",
        "2. It does not know about your local Dynamo endpoint setup",
        "Your setup uses:",
        "local Dynamo frontend",
        "local SGLang worker",
        "OpenAI-compatible local endpoint",
        "nvext.agent_hints",
        "That is not the default expected environment for the upstream examples.",
    ],
    top=1.4,
    font_size=16,
)

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "What Upstream Does Not Know By Default (Cont.)", "Added verbatim from the AgentBench setup explanation")
add_bullets(
    slide,
    [
        "3. It does not give you your checkpoint observability layer",
        "You wanted explicit checkpoints:",
        "1",
        "3",
        "4",
        "5",
        "with saved payloads in:",
        "checkpoints.json",
        "That is your experiment instrumentation, not generic upstream behavior.",
        "",
        "4. It does not save the exact artifacts you want",
        "You wanted things like:",
        "result.json",
        "plan.json",
        "step_results.json",
        "final_summary.txt",
        "checkpoints.json",
        "Again, that is benchmark/testbed wrapper behavior.",
    ],
    top=1.4,
    font_size=16,
)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Where The Setup Hooks Live", "Exact wrapper and app-layer hook points for this setup")
add_bullets(
    slide,
    [
        "1. SWE-bench task ingestion",
        "This is where your setup teaches the harness how to read ScaleAI/SWE-bench_Pro or local task files.",
        "agentbench/deepagents_swebench_single_host.py: load_swebench_task(...)",
        "main() calls load_swebench_task(...).",
        "",
        "2. SWE-bench repo checkout adaptation",
        "This is where your setup turns SWE-bench task metadata into a real GitHub repo workspace.",
        "agentbench/deepagents_swebench_single_host.py: infer_swebench_repo_url(task)",
        "agentbench/deepagents_swebench_single_host.py: infer_swebench_base_commit(task)",
        "agentbench/deepagents_swebench_single_host.py: ensure_shared_repo_checkout(...)",
        "agentbench/deepagents_swebench_single_host.py: prepare_workspace(...)",
        "main() decides whether to auto-materialize a SWE-bench repo.",
        "",
        "3. Hand-off from wrapper to Deep Agents app",
        "agentbench/deepagents_swebench_single_host.py: workflow = run_task_workflow(...)",
    ],
    top=1.35,
    font_size=15,
)

# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Where The Setup Hooks Live (Cont.)", "Exact wrapper and app-layer hook points for this setup")
add_bullets(
    slide,
    [
        "4. Deep Agents runtime source adaptation",
        "This is where your app prefers the cloned GitHub Deep Agents code instead of only the installed package.",
        "agentbench/deepagents_app/src/agent.py: UPSTREAM_ROOT",
        "agentbench/deepagents_app/src/agent.py: CLONED_DEEPAGENTS_LIB_ROOT",
        "agentbench/deepagents_app/src/agent.py: prepends cloned Deep Agents path to sys.path",
        "agentbench/deepagents_app/src/agent.py: from deepagents import create_deep_agent",
        "",
        "5. Dynamo frontend adaptation",
        "This is where your app tells Deep Agents to talk to your local Dynamo endpoint.",
        "agentbench/deepagents_app/src/agent.py: frontend_base_url(...)",
        "agentbench/deepagents_app/src/agent.py: build_dynamo_chat_model(...)",
        "This is the exact place where Deep Agents is adapted to your Dynamo frontend.",
        "",
        "6. SGLang connection",
        "There is no direct SGLang call inside agentbench/.",
        "Inside agentbench, the SGLang hook is really: send everything to the Dynamo frontend URL.",
    ],
    top=1.35,
    font_size=15,
)

# Slide 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Where The Setup Hooks Live (Cont.)", "Exact wrapper and app-layer hook points for this setup")
add_bullets(
    slide,
    [
        "7. Deep Agents harness creation",
        "This is where your app builds the actual agent object.",
        "agentbench/deepagents_app/src/agent.py: build_coding_agent(...)",
        "agentbench/deepagents_app/src/agent.py: create_deep_agent(...)",
        "",
        "8. Planning / step / synthesis adaptation",
        "This is where your wrapper-app defines the explicit multi-step workflow.",
        "agentbench/deepagents_app/src/agent.py: generate_decomposition_plan(...)",
        "agentbench/deepagents_app/src/agent.py: execute_plan_steps(...)",
        "agentbench/deepagents_app/src/agent.py: synthesize_final_summary(...)",
        "agentbench/deepagents_app/src/agent.py: run_task_workflow(...)",
        "",
        "9. Checkpoint observability hooks",
        "deepagents_swebench_single_host.py: checkpoint 1",
        "agent.py: checkpoint 3",
        "agent.py: checkpoint 4",
        "agent.py: checkpoint 5",
    ],
    top=1.35,
    font_size=15,
)

# Slide 11
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Where The Setup Hooks Live (Cont.)", "Exact wrapper and app-layer hook points for this setup")
add_bullets(
    slide,
    [
        "10. Saved artifacts hook",
        "This is where the run outputs are written into your agentbench/results/... format.",
        "agentbench/deepagents_swebench_single_host.py: save_result(...)",
        "writes plan.json",
        "writes step_results.json",
        "writes final_summary.txt",
        "captures workspace artifacts",
        "builds final result.json payload",
        "",
        "So the full adaptation chain is:",
        "SWE-bench task load -> repo/commit checkout -> handoff to Deep Agents app -> Deep Agents via cloned GitHub source -> Dynamo frontend wiring -> Deep Agents workflow -> saved artifacts",
    ],
    top=1.35,
    font_size=15,
)

# Slide 12
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Concrete Walkthrough", "What happens during one AgentBench run")
add_bullets(
    slide,
    [
        "Use this concrete example task from AgentBench:",
        "A Python CLI has a --dry-run flag. Users say the command still creates a temporary working directory during dry run.",
        "",
        "And imagine you run:",
        "cd ~/kv_cache_offloading",
        "bash agentbench/run_upstream_deploy_coding_agent_single_host.sh --dataset ScaleAI/SWE-bench_Pro --split test --index 0",
        "",
        "Here is what happens, stage by stage.",
    ],
    top=1.45,
    font_size=16,
)

# Slide 13
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Concrete Walkthrough (Cont.)", "Stages 1 to 3")
add_bullets(
    slide,
    [
        "1. Task load",
        "The wrapper script agentbench/deepagents_swebench_single_host.py loads one task.",
        "What it gets:",
        "instance_id",
        "repo",
        "problem_statement",
        "maybe base_commit",
        "other SWE-bench fields",
        "At this point it is just a Python dictionary.",
        "",
        "Concrete example",
        "It might load:",
        'repo = "owner/project"',
        'base_commit = "abc123"',
        'problem_statement = "...dry-run still creates temp directory..."',
        "This is checkpoint 1.",
        "",
        "2. Repo materialization",
        "The wrapper looks at the task metadata and turns it into a real repo workspace.",
    ],
    top=1.3,
    font_size=15,
)

# Slide 14
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Concrete Walkthrough (Cont.)", "Stages 2 to 4")
add_bullets(
    slide,
    [
        "What it does:",
        'converts owner/project into https://github.com/owner/project.git',
        "keeps a shared checkout under agentbench/repos/owner__project/",
        "checks out abc123 there if available",
        "So the task is no longer just text.",
        "Now it also has a real repo on disk.",
        "",
        "3. Handoff into the Deep Agents app",
        "The wrapper then calls:",
        "run_task_workflow(...)",
        "in agentbench/deepagents_app/src/agent.py",
        "This is the boundary between:",
        "outer benchmark wrapper",
        "inner agent app",
        "",
        "4. Task prompt construction",
        "Inside the app, the task dictionary is turned into one large task prompt.",
    ],
    top=1.3,
    font_size=15,
)

# Slide 15
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Concrete Walkthrough (Cont.)", "Stages 4 to 6")
add_bullets(
    slide,
    [
        "That prompt includes things like:",
        "task metadata",
        "repo name",
        "problem statement",
        "workspace path if present",
        "So the raw SWE-bench row becomes model-readable prompt text.",
        "",
        "5. Planning stage",
        "The app creates a planning request.",
        "What it asks:",
        "break this task into steps",
        "return JSON",
        "keep it under a step limit",
        "prefer task-specific steps",
        "For the dry-run example, it tries to get steps like:",
        "inspect argument parsing for --dry-run",
        "inspect dispatch flow",
        "inspect temp-dir helper",
        "propose safe fix strategy",
    ],
    top=1.3,
    font_size=15,
)

# Slide 16
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Concrete Walkthrough (Cont.)", "Stages 5 to 8")
add_bullets(
    slide,
    [
        "This request is sent to:",
        "local Dynamo frontend",
        "This is checkpoint 3.",
        "",
        "6. Dynamo frontend",
        "The app does not call SGLang directly.",
        "Instead it builds a ChatOpenAI(...) client with:",
        "local base_url",
        "nvext.agent_hints",
        "So the request goes to:",
        "http://127.0.0.1:8000/v1/chat/completions",
        "Dynamo receives the planning request and forwards it to the local SGLang worker.",
        "",
        "7. SGLang worker",
        "SGLang runs the actual model.",
        "It returns the planning response back through Dynamo to AgentBench.",
    ],
    top=1.3,
    font_size=15,
)

# Slide 17
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Where The Hints Come From", "Added from the AgentBench setup explanation")
add_bullets(
    slide,
    [
        "The hints come from your own AgentBench code, not from Dynamo and not from Deep Agents automatically.",
        "",
        "Short answer",
        "The values come from three places:",
        "1. default hints hardcoded in AgentBench",
        "2. optional base hints passed in from the runner",
        "3. per-phase overrides added by the app for planning, step execution, and synthesis",
        "",
        "Where they start",
        "In the runner:",
        "agentbench/deepagents_swebench_single_host.py: DEFAULT_HINTS = {...}",
        "That includes values like:",
        "priority",
        "reuse_likelihood",
        "agent_phase",
        "latency_sensitivity",
        "program_id",
        "context_type",
        "expected_output_tokens",
    ],
    top=1.25,
    font_size=15,
)

# Slide 18
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Where The Hints Come From (Cont.)", "Added from the AgentBench setup explanation")
add_bullets(
    slide,
    [
        "The CLI also exposes:",
        "agentbench/deepagents_swebench_single_host.py: --hint-json",
        "Then in main() it does:",
        "base_hints = json.loads(args.hint_json)",
        "So that is the first source.",
        "",
        "Where they are resolved in the app",
        "In the app layer:",
        "agentbench/deepagents_app/src/agent.py: DEFAULT_DYNAMO_HINTS = {...}",
        "Then:",
        "build_phase_hints(...) does:",
        "1. start from DEFAULT_DYNAMO_HINTS",
        "2. merge in base_hints from the runner",
        "3. overwrite agent_phase with the current phase",
        "So the app decides the final hint payload for each request.",
    ],
    top=1.25,
    font_size=15,
)

# Slide 19
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Where The Hints Come From (Cont.)", "Added from the AgentBench setup explanation")
add_bullets(
    slide,
    [
        "How it knows what values to use",
        "Right now, it knows because you hardcoded the policy.",
        "Example:",
        "planning should be a little less latency-sensitive",
        "planning should expect smaller output",
        "step execution should expect more output",
        "synthesis should expect more output",
        "Those choices are written directly in the code.",
        "",
        "Planning",
        "planning_hints = build_phase_hints(base_hints, phase=\"planning\")",
        "planning_hints[\"latency_sensitivity\"] = 0.4",
        "planning_hints[\"expected_output_tokens\"] = 512",
        "",
        "Step execution",
        "step_hints = build_phase_hints(base_hints, phase=f\"step_{idx}_execution\")",
        "step_hints[\"expected_output_tokens\"] = 768",
    ],
    top=1.25,
    font_size=15,
)

# Slide 20
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Where The Hints Come From (Cont.)", "Added from the AgentBench setup explanation")
add_bullets(
    slide,
    [
        "Synthesis",
        "synthesis_hints = build_phase_hints(base_hints, phase=\"synthesis\")",
        "synthesis_hints[\"expected_output_tokens\"] = 768",
        "",
        "So the answer is:",
        "it does not discover these values automatically",
        "it uses your chosen heuristic values",
        "",
        "How the hints get attached to the request",
        "That happens in:",
        "extra_body = {\"nvext\": {\"agent_hints\": payload}}",
        "and then in the ChatOpenAI(...) client.",
        "So every model call to Dynamo carries:",
        "nvext.agent_hints = <resolved payload>",
        "",
        "One-line answer",
        "The hints are currently hand-authored heuristics in AgentBench, merged and adjusted per phase before being attached to each Dynamo request as nvext.agent_hints.",
    ],
    top=1.25,
    font_size=15,
)

# Slide 21
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Hints In Very Simple Terms", "Direct answer to the SWE-bench dataset question")
add_bullets(
    slide,
    [
        "Are the hints gotten from the SWE dataset?",
        "No.",
        "",
        "Very simple explanation",
        "The SWE-bench dataset gives AgentBench the task.",
        "AgentBench itself makes up the hints.",
        "Those hints are small routing/serving instructions that AgentBench attaches before sending the request to Dynamo.",
        "",
        "So:",
        "SWE-bench gives the problem.",
        "AgentBench gives the hints.",
        "Dynamo receives the request plus those hints.",
        "",
        "Shortest version",
        "No, the hints are not coming from SWE-bench. They are defined in your AgentBench code.",
    ],
    top=1.35,
    font_size=16,
)

# Slide 22
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Concrete Walkthrough (Cont.)", "Stages 7 to 9")
add_bullets(
    slide,
    [
        "If the model behaves well, AgentBench gets a clean JSON plan.",
        "If not, AgentBench tries to recover a plan with its fallback parser.",
        "",
        "8. Step execution stage",
        "Now AgentBench loops over the plan steps.",
        "For each step, it creates a new request like:",
        "here is the original task",
        "here is the approved plan",
        "here is the current step",
        "here are completed step summaries so far",
        "",
        "Example step:",
        "Inspect argument parsing for --dry-run",
        "",
        "That step request again goes:",
        "Deep Agents app -> Dynamo -> SGLang",
        "Each of these is checkpoint 4.",
    ],
    top=1.3,
    font_size=15,
)

# Slide 23
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Concrete Walkthrough (Cont.)", "Stages 8 to 10")
add_bullets(
    slide,
    [
        "If a repo workspace exists, the agent runs with that repo as its working directory.",
        "",
        "9. Synthesis stage",
        "After all steps finish, AgentBench creates one final synthesis prompt.",
        "It includes:",
        "the original task",
        "the plan used",
        "all step summaries",
        "It asks for:",
        "overall diagnosis",
        "proposed fix strategy",
        "files/code areas that matter",
        "validation still needed",
        "what changed in the workspace",
        "This final request again goes:",
        "Deep Agents app -> Dynamo -> SGLang",
        "This is checkpoint 5.",
    ],
    top=1.3,
    font_size=15,
)

# Slide 24
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Concrete Walkthrough (Cont.)", "Final save stage")
add_bullets(
    slide,
    [
        "10. Saved artifacts",
        "After the final response comes back, the wrapper saves everything under:",
        "agentbench/results/<run>/",
        "Typical files:",
        "result.json",
        "plan.json",
        "step_results.json",
        "final_summary.txt",
        "checkpoints.json",
        "If a repo workspace was used, it may also save:",
        "workspace.patch",
        "git_status.txt",
        "git_diff_stat.txt",
        "",
        "End-to-end in one line",
        "SWE-bench task row -> wrapper loads task -> wrapper materializes repo -> app builds prompt -> planning request -> Dynamo -> SGLang -> step requests -> Dynamo -> SGLang -> synthesis request -> Dynamo -> SGLang -> artifacts saved",
    ],
    top=1.3,
    font_size=15,
)

# Slide 25
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Concrete Walkthrough (Cont.)", "Layer responsibilities")
add_bullets(
    slide,
    [
        "What each layer is responsible for",
        "Wrapper:",
        "task loading",
        "repo checkout",
        "run directory",
        "artifact saving",
        "Deep Agents app:",
        "prompt building",
        "planning",
        "step execution",
        "synthesis",
        "Dynamo hint wiring",
        "Dynamo:",
        "request entrypoint",
        "forwards model calls",
        "SGLang:",
        "actual model execution",
    ],
    top=1.3,
    font_size=15,
)

# Slide 26
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Experiment Modes", "Main ways to run the upstream deploy-coding-agent variant")
cmd_text = (
    "Sample task\n"
    "bash agentbench/run_upstream_deploy_coding_agent_single_host.sh\n\n"
    "SWE-bench Pro input\n"
    "bash agentbench/run_upstream_deploy_coding_agent_single_host.sh \\\n"
    "  --dataset ScaleAI/SWE-bench_Pro \\\n"
    "  --split test \\\n"
    "  --index 0\n\n"
    "Stronger-behavior sample run used in this deck\n"
    "bash agentbench/run_upstream_deploy_coding_agent_stronger_behavior_single_host.sh"
)
box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(12.0), Inches(4.6))
box.fill.solid()
box.fill.fore_color.rgb = NAVY
box.line.color.rgb = NAVY
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = cmd_text
r.font.name = "Courier New"
r.font.size = Pt(17)
r.font.color.rgb = WHITE
add_bullets(
    slide,
    [
        "This deck uses the stronger-behavior run so the checkpoint slides are filled with real outputs.",
    ],
    top=6.3,
    font_size=17,
)

# Slide 27
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Checkpoint Map", "Ordered observability points in the upstream pipeline")
add_box(slide, 0.8, 1.55, 2.5, 0.8, "1\nTask loaded before\nDeep Agents harness", SOFT_BLUE, LIGHT, 17, True)
add_box(slide, 3.6, 1.55, 2.5, 0.8, "3\nPlanning request leaving\nDeep Agents harness", SOFT_GREEN, LIGHT, 17, True)
add_box(slide, 6.4, 1.55, 2.5, 0.8, "4\nStep execution request\nleaving harness", SOFT_ORANGE, LIGHT, 17, True)
add_box(slide, 9.2, 1.55, 2.5, 0.8, "5\nFinal synthesis request\nleaving harness", SOFT_PINK, LIGHT, 17, True)
for x in [3.2, 6.0, 8.8]:
    add_arrow(slide, x, 1.75)
add_bullets(
    slide,
    [
        "1: measured at the point just before the loaded task enters the Deep Agents harness.",
        "3: measured at the point where the planning stage is about to send its first request to the Dynamo frontend.",
        "4: measured at the point where the execution stage is about to send one step-level request to the Dynamo frontend.",
        "5: measured at the point where the synthesis stage is about to send the final combined request to the Dynamo frontend.",
    ],
    top=3.0,
    font_size=18,
)

# Slide 28
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Observed Run Configuration", "Actual run used to fill the checkpoint slides")
add_box(slide, 0.8, 1.5, 4.0, 2.3, "Run ID\nsample__stronger_behavior__001_20260506_155649\n\nTask source\njson:agentbench/sample_task_stronger.json\n\nInstance ID\nsample__stronger_behavior__001", WHITE, LIGHT, 16, True)
add_box(slide, 5.1, 1.5, 3.1, 2.3, "App variant\nupstream_deploy_coding_agent\n\nModel\nQwen/Qwen2.5-0.5B\n\nRuntime source\npython_environment", SOFT_BLUE, LIGHT, 16, True)
add_box(slide, 8.5, 1.5, 3.8, 2.3, "Workspace\nnone provided\n\nPlan step count\n4\n\nCheckpoint file\ncheckpoints.json", SOFT_GREEN, LIGHT, 16, True)
add_bullets(
    slide,
    [
        f"Problem statement: {shorten(task['problem_statement'], 180)}",
        "This run was chosen because it shows Checkpoint 1, Checkpoint 3, four separate Checkpoint 4 events, and Checkpoint 5.",
    ],
    top=4.25,
    font_size=18,
)

# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Checkpoint 1", "Task ingestion stage")
add_sections_text(
    slide,
    [
        ("Represents", ["The loaded task just before it enters the Deep Agents harness."]),
        (
            "What changed here",
            [
                "The external task source was normalized into one task record.",
                "No planning prompt exists yet.",
                "No phase-specific hints are attached yet.",
            ],
        ),
    ],
    body_size=15,
    section_title_size=18,
)

# Slide 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Checkpoint 1", "Payload at this stage from checkpoints.json")
add_preformatted_text(slide, raw_json_lines(checkpoint_1, max_string=220, max_lines=28), font_size=12)

# Slide 11
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Checkpoint 3", "Planning request stage")
add_sections_text(
    slide,
    [
        (
            "Represents",
            [
                "The first outbound planning request, just before it goes to the Dynamo frontend.",
            ],
        ),
        (
            "What changed here",
            [
                "The task record became a planning request.",
                "Planning-stage hints were attached.",
                "The request became frontend-bound.",
            ],
        ),
    ],
    body_size=15,
    section_title_size=18,
)

# Slide 12
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Checkpoint 3", "Payload at this stage from checkpoints.json")
add_preformatted_text(slide, raw_json_lines(checkpoint_3, max_string=160, max_lines=30), font_size=10.5)

# Slide 13
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Checkpoint 4", "Per-step execution request stage")
add_sections_text(
    slide,
    [
        (
            "Represents",
            [
                "One outbound request for one current step after a decomposition plan exists.",
                f"In this run, Checkpoint 4 appears {len(checkpoint_4s)} times, once per planned step.",
                "The next slides show all four saved step-level payloads from checkpoints.json.",
            ],
        ),
        (
            "What changed here",
            [
                "The single planning output became multiple step-level requests.",
                "Each request gained step_index and step_title.",
                "The hint phase changed from planning to step_n_execution.",
            ],
        ),
    ],
    body_size=15,
    section_title_size=18,
)

for idx, checkpoint_4 in enumerate(checkpoint_4s, start=1):
    # Additional checkpoint 4 payload slides
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Checkpoint 4",
        f"Payload {idx} of {len(checkpoint_4s)} from checkpoints.json: {shorten(checkpoint_4['step_title'], 70)}",
    )
    add_preformatted_text(slide, raw_json_lines(checkpoint_4, max_string=160, max_lines=30), font_size=10.5)

# After checkpoint 4 payload slides
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Checkpoint 5", "Final synthesis request stage")
add_sections_text(
    slide,
    [
        (
            "Represents",
            [
                "The final outbound request after all step-level requests have completed.",
                "This is the merge stage that turns multiple step results into one final request.",
            ],
        ),
        (
            "What changed here",
            [
                "Per-step results were aggregated into one request.",
                "The request became task-level again instead of step-level.",
                "The hint phase changed from step_n_execution to synthesis.",
            ],
        ),
    ],
    body_size=15,
    section_title_size=18,
)

# Next slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Checkpoint 5", "Payload at this stage from checkpoints.json")
add_preformatted_text(slide, raw_json_lines(checkpoint_5, max_string=160, max_lines=30), font_size=10.5)

# Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Payload Evolution Summary", "How the structure changes from checkpoint to checkpoint")
add_box(slide, 0.55, 1.55, 2.7, 1.0, "1\nTask record", SOFT_BLUE, LIGHT, 20, True)
add_box(slide, 3.45, 1.55, 2.9, 1.0, "3\nPlanning request", SOFT_GREEN, LIGHT, 20, True)
add_box(slide, 6.6, 1.55, 2.9, 1.0, "4\nStep request(s)", SOFT_ORANGE, LIGHT, 20, True)
add_box(slide, 9.75, 1.55, 2.8, 1.0, "5\nSynthesis request", SOFT_PINK, LIGHT, 20, True)
for x in [3.1, 6.25, 9.4]:
    add_arrow(slide, x, 1.82, width=0.45)
add_sections_text(
    slide,
    [
        (
            "Structure change by stage",
            [
                "Checkpoint 1: task record only.",
                "Checkpoint 3: task record + planning wrapper + planning hints.",
                "Checkpoint 4: current step + step wrapper + step-specific hints.",
                "Checkpoint 5: aggregated step results + synthesis wrapper + synthesis hints.",
            ],
        ),
    ],
    body_size=16,
    section_title_size=18,
)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Wrote {OUT}")
